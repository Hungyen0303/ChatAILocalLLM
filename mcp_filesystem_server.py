#!/usr/bin/env python3
"""
MCP Filesystem Server
Quản lý filesystem cho ứng dụng Chat AI Local LLM
Chức năng: Tìm kiếm, index và phân loại file văn bản (PDF, Word, PPT)
"""

import asyncio
import json
import os
import logging
from typing import Dict, List, Any, Optional
from pathlib import Path
import mcp.types as types
from mcp.server import Server
from mcp.server.models import InitializationOptions
import mcp.server.stdio
from pydantic import BaseModel
from llm_utils import ask_llm_yesno

# Import cấu hình đơn giản
from config import SUPPORTED_EXTENSIONS, CONTENT_PREVIEW_LIMIT, CATEGORY_KEYWORDS

# Import thư viện xử lý file
try:
    import PyPDF2
    import docx
    from pptx import Presentation
except ImportError:
    print("Cài đặt thêm: pip install PyPDF2 python-docx python-pptx")

# Cấu hình logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class FileMetadata(BaseModel):
    """Metadata của file"""
    filename: str
    filepath: str
    file_type: str
    size: int
    content_preview: str
    label: str = "Chưa phân loại"
    created_time: float
    modified_time: float

class FileIndexer:
    """Class quản lý index file"""
    MCP_CLOUD_API_URL = "http://localhost:8000/upload-metadata" 
    
    def __init__(self, base_path: str = "."):
        self.base_path = Path(base_path)
        self.file_index: Dict[str, FileMetadata] = {}
        self.supported_extensions = set(SUPPORTED_EXTENSIONS)
    
    def extract_text_from_pdf(self, filepath: Path) -> str:
        """Trích xuất text từ file PDF"""
        try:
            with open(filepath, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text[:CONTENT_PREVIEW_LIMIT]
        except Exception as e:
            logger.error(f"Lỗi đọc PDF {filepath}: {e}")
            return ""
    
    def extract_text_from_docx(self, filepath: Path) -> str:
        """Trích xuất text từ file Word"""
        try:
            doc = docx.Document(filepath)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text[:CONTENT_PREVIEW_LIMIT]
        except Exception as e:
            logger.error(f"Lỗi đọc DOCX {filepath}: {e}")
            return ""
    
    def extract_text_from_pptx(self, filepath: Path) -> str:
        """Trích xuất text từ file PowerPoint"""
        try:
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text[:CONTENT_PREVIEW_LIMIT]
        except Exception as e:
            logger.error(f"Lỗi đọc PPTX {filepath}: {e}")
            return ""
    
    def extract_text_from_txt(self, filepath: Path) -> str:
        """Trích xuất text từ file TXT"""
        try:
            with open(filepath, 'r', encoding='utf-8') as file:
                text = file.read()
                return text[:CONTENT_PREVIEW_LIMIT]
        except Exception as e:
            logger.error(f"Lỗi đọc TXT {filepath}: {e}")
            return ""
    
    def extract_content(self, filepath: Path) -> str:
        """Trích xuất nội dung từ file dựa trên extension"""
        extension = filepath.suffix.lower()
        
        if extension == '.pdf':
            return self.extract_text_from_pdf(filepath)
        elif extension in ['.docx', '.doc']:
            return self.extract_text_from_docx(filepath)
        elif extension in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx(filepath)
        elif extension == '.txt':
            return self.extract_text_from_txt(filepath)
        else:
            return ""
    
    
    def scan_directory(self, directory: Path = None) -> List[FileMetadata]:
        """Quét thư mục và tạo index file"""
        if directory is None:
            directory = self.base_path
        
        files_found = []
        
        for filepath in directory.rglob("*"):
            if filepath.is_file() and filepath.suffix.lower() in self.supported_extensions:
                try:
                    # Lấy thông tin file
                    stat = filepath.stat()
                    content = self.extract_content(filepath)
                    label = "Chưa phân loại"
                    
                    # Tạo metadata
                    metadata = FileMetadata(
                        filename=filepath.name,
                        filepath=str(filepath.absolute()),
                        file_type=filepath.suffix.lower(),
                        size=stat.st_size,
                        content_preview=content,
                        label=label,
                        created_time=stat.st_ctime,
                        modified_time=stat.st_mtime
                    )
                    self.file_index[str(filepath.absolute())] = metadata
                    files_found.append(metadata)
                    
                    logger.info(f"Indexed: {filepath.name} -> {label}")
                    
                except Exception as e:
                    logger.error(f"Lỗi index file {filepath}: {e}")
        
        return files_found
    
    def search_files(self, query: str) -> List[FileMetadata]:
        """Tìm kiếm file theo query"""
        query_lower = query.lower()
        results = []
        
        for metadata in self.file_index.values():
            # Tìm trong tên file và nội dung
            if (query_lower in metadata.filename.lower() or 
                query_lower in metadata.content_preview.lower()):
                results.append(metadata)
        return results
    
    def get_files_by_category(self, category: str) -> List[FileMetadata]:
        """Lấy file theo nhóm phân loại"""
        results = []
        for metadata in self.file_index.values():
            if category.lower() in metadata.label.lower():
                results.append(metadata)
        return results

    def extract_full_content(self, filepath: Path) -> str:
        extension = filepath.suffix.lower()
        try:
            if extension == '.pdf':
                with open(filepath, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    return "\n".join([page.extract_text() or "" for page in reader.pages])
            elif extension in ['.docx', '.doc']:
                doc = docx.Document(filepath)
                return "\n".join([p.text for p in doc.paragraphs])
            elif extension in ['.pptx', '.ppt']:
                prs = Presentation(filepath)
                return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
            elif extension == '.txt':
                with open(filepath, 'r', encoding='utf-8') as file:
                    return file.read()
            else:
                return ""
        except Exception as e:
            logger.error(f"Lỗi đọc toàn bộ nội dung {filepath}: {e}")
            return ""

    def classify_files_by_topic(self, topic: str):
        """Cập nhật label cho từng file nếu liên quan chủ đề"""
        for metadata in self.file_index.values():
            full_content = self.extract_full_content(Path(metadata.filepath))
            if len(full_content) > 4000:
                full_content = full_content[:3000] + '\n...\n' + full_content[-1000:]
            if not full_content.strip():
                continue
            if ask_llm_yesno(full_content, topic):
                metadata.label = f"{topic}"
        
# Khởi tạo file indexer
file_indexer = FileIndexer()

# Tạo MCP Server
server = Server("filesystem-manager")

@server.list_resources()
async def handle_list_resources() -> list[types.Resource]:
    """Liệt kê các resource có sẵn"""
    return [
        types.Resource(
            uri="filesystem://index",
            name="File Index",
            description="Index của tất cả file trong hệ thống",
            mimeType="application/json",
        ),
        types.Resource(
            uri="filesystem://search",
            name="File Search",
            description="Tìm kiếm file theo từ khóa",
            mimeType="application/json",
        ),
    ]

@server.read_resource()
async def handle_read_resource(uri: str) -> str:
    """Đọc resource theo URI"""
    if uri == "filesystem://index":
        # Trả về danh sách tất cả file đã index
        files = list(file_indexer.file_index.values())
        return json.dumps([file.dict() for file in files], indent=2, ensure_ascii=False)
    elif uri.startswith("filesystem://search?q="):
        # Tìm kiếm file
        query = uri.split("q=")[1]
        results = file_indexer.search_files(query)
        return json.dumps([file.dict() for file in results], indent=2, ensure_ascii=False)
    else:
        raise ValueError(f"Unknown resource: {uri}")

@server.list_tools()
async def handle_list_tools() -> list[types.Tool]:
    """Liệt kê các tool có sẵn"""
    return [
        types.Tool(
            name="scan_directory",
            description="Quét thư mục và tạo index file",
            inputSchema={
                "type": "object",
                "properties": {
                    "directory": {
                        "type": "string",
                        "description": "Đường dẫn thư mục cần quét (để trống = thư mục hiện tại)"
                    }
                }
            },
        ),
        types.Tool(
            name="search_files",
            description="Tìm kiếm file theo từ khóa",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Từ khóa tìm kiếm"
                    }
                },
                "required": ["query"]
            },
        ),
        types.Tool(
            name="get_file_info",
            description="Lấy thông tin chi tiết của file",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Đường dẫn file"
                    }
                },
                "required": ["filepath"]
            },
        ),
        types.Tool(
            name="export_metadata",
            description="Xuất metadata để gửi qua MCP Cloud",
            inputSchema={
                "type": "object",
                "properties": {
                    "format": {
                        "type": "string",
                        "description": "Định dạng xuất (json, csv)",
                        "default": "json"
                    }
                }
            },
        ),
        types.Tool(
            name="classify_files_by_topic",
            description="Phân loại file liên quan đến chủ đề/từ khóa tự nhiên",
            inputSchema={
                "type": "object",
                "properties": {
                    "topic": {
                        "type": "string",
                        "description": "Chủ đề hoặc từ khóa tự nhiên"
                    }
                },
                "required": ["topic"]
            },
        ),
    ]

@server.call_tool()
async def handle_call_tool(name: str, arguments: dict) -> list[types.TextContent]:
    """Xử lý tool calls"""
    
    if name == "scan_directory":
        directory = arguments.get("directory", ".")
        try:
            files = file_indexer.scan_directory(Path(directory))
            result = {
                "message": f"Đã quét và index {len(files)} file",
                "files": [
                    {
                        "filename": f.filename,
                        "label": f.label,
                        "size": f.size,
                        "content_preview": f.content_preview[:200] + "..." if len(f.content_preview) > 200 else f.content_preview,
                        "type": f.file_type
                    } for f in files
                ]
            }
            return [types.TextContent(type="text", text=json.dumps(result, indent=2, ensure_ascii=False))]
        except Exception as e:
            return [types.TextContent(type="text", text=f"Lỗi quét thư mục: {e}")]
    
    elif name == "search_files":
        query = arguments["query"]
        try:
            results = file_indexer.search_files(query)
            result = {
                "query": query,
                "found": len(results),
                "files": [
                    {
                        "filename": f.filename,
                        "filepath": f.filepath,
                        "label": f.label,
                        "content_preview": f.content_preview[:200] + "..." if len(f.content_preview) > 200 else f.content_preview
                    } for f in results
                ]
            }
            return [types.TextContent(type="text", text=json.dumps(result, indent=2, ensure_ascii=False))]
        except Exception as e:
            return [types.TextContent(type="text", text=f"Lỗi tìm kiếm: {e}")]
    
    elif name == "get_file_info":
        filepath = arguments["filepath"]
        print (f"filepath {filepath}")
        try:
            metadata = file_indexer.file_index.get(filepath)
            if metadata:
                return [types.TextContent(type="text", text=json.dumps(metadata.dict(), indent=2, ensure_ascii=False))]
            else:
                return [types.TextContent(type="text", text=f"Không tìm thấy file: {filepath}")]
        except Exception as e:
            return [types.TextContent(type="text", text=f"Lỗi lấy thông tin file: {e}")]
    
    elif name == "export_metadata":
        format_type = arguments.get("format", "json")
        try:
            files = list(file_indexer.file_index.values())
            

            if format_type == "json":
                metadata_for_cloud = []
                success_count = 0
                error_count = 0
                error_files = []

                for f in files:
                    metadata = {
                        "filename": f.filename,
                        "label": f.label,
                        "content": f.content_preview[:500],
                        "file_type": f.file_type,
                        "size": f.size
                    }
                    metadata_for_cloud.append(metadata)

                    # Gửi metadata lên MCP Cloud
                    try:
                        resp = requests.post(MCP_CLOUD_API_URL, json=metadata)
                        if resp.status_code == 200:
                            success_count += 1
                        else:
                            error_count += 1
                            error_files.append(f.filename)
                    except Exception as e:
                        error_count += 1
                        error_files.append(f.filename)
                
                logger.info(f"Đã gửi {success_count} metadata thành công, {error_count} lỗi")
                if error_files:
                    logger.error(f"Các file gặp lỗi: {', '.join(error_files)}")

                return [types.TextContent(type="text", text=json.dumps(metadata_for_cloud, indent=2, ensure_ascii=False))]
            else:
                return [types.TextContent(type="text", text="Chỉ hỗ trợ định dạng JSON")]
        except Exception as e:
            return [types.TextContent(type="text", text=f"Lỗi xuất metadata: {e}")]
    
    elif name == "classify_files_by_topic":
        topic = arguments.get("topic", "")
        try:
            # Gán label cho từng file 
            file_indexer.classify_files_by_topic(topic)
            # Gom nhóm các file đã được gán label
            results = file_indexer.get_files_by_category(f"Liên quan: {topic}")
            result = {
                "topic": topic,
                "count": len(results),
                "files": [
                    {
                        "filename": f.filename,
                        "label": f.label,
                        "filepath": f.filepath
                    } for f in results
                ]
            }
            return [types.TextContent(type="text", text=json.dumps(result, indent=2, ensure_ascii=False))]
        except Exception as e:
            return [types.TextContent(type="text", text=f"Lỗi phân loại theo chủ đề: {e}")]
    
    else:
        return [types.TextContent(type="text", text=f"Tool không được hỗ trợ: {name}")]

async def main():
    """Chạy MCP server"""
    # Quét thư mục ban đầu
    logger.info("Khởi động MCP Filesystem Server...")
    logger.info("Quét thư mục ban đầu...")
    try:
        initial_files = file_indexer.scan_directory()
        logger.info(f"Đã index {len(initial_files)} file ban đầu")
    except Exception as e:
        logger.error(f"Lỗi quét thư mục ban đầu: {e}")
    
    # Chạy server
    async with mcp.server.stdio.stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="filesystem-manager",
                server_version="1.0.0",
                capabilities=server.get_capabilities(
                    notification_options=None,
                    experimental_capabilities=None,
                ),
            ),
        )

if __name__ == "__main__":
    asyncio.run(main()) 